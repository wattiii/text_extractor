"""
extract_documents_to_csv_xlsx.py

Purpose:
- Walk through a directory with subfolders
- Extract text from .docx, .pptx, .xlsx, .pdf files
- Chunk by sentence
- Capture metadata (file path, type, page/slide/sheet, comments)
- Save both CSV and XLSX for easy consumption (Power BI, Excel)
"""

import os
import pandas as pd
import docx
import pptx
import openpyxl
import pdfplumber
import pytesseract
from PIL import Image
from pdf2image import convert_from_path
import re
from tqdm import tqdm
import nltk

# Ensure NLTK sentence tokenizer is downloaded
nltk.download('punkt')
from nltk.tokenize import sent_tokenize

# -------------------------------
# Helper functions
# -------------------------------

def extract_docx(file_path):
    doc = docx.Document(file_path)
    text = []
    for para in doc.paragraphs:
        if para.text.strip():
            text.append(para.text.strip())
    # Optionally capture comments
    comments = []
    if doc.part.element.xpath("//w:comment"):
        comments = [c.text for c in doc.part.element.xpath("//w:comment")]
    return text, comments

def extract_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        text.append((i + 1, " ".join(slide_text)))
    return text  # list of tuples: (slide_number, slide_text)

def extract_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path, data_only=True)
    rows = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
            row_text = " | ".join([str(cell) for cell in row if cell is not None])
            if row_text.strip():
                rows.append((sheet_name, row_idx, row_text))
    return rows  # list of tuples: (sheet_name, row_number, text)

def extract_pdf(file_path, use_ocr=True):
    text_items = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_items.append((i+1, page_text))
                elif use_ocr:
                    # Fallback OCR if page is image-heavy
                    pil_images = convert_from_path(file_path, first_page=i+1, last_page=i+1)
                    for img in pil_images:
                        ocr_text = pytesseract.image_to_string(img)
                        if ocr_text.strip():
                            text_items.append((i+1, ocr_text))
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
    return text_items  # list of tuples: (page_number, text)

def chunk_sentences(text):
    return sent_tokenize(text)

# -------------------------------
# Main directory walk
# -------------------------------
def process_directory(root_dir):
    records = []
    for dirpath, _, filenames in os.walk(root_dir):
        for filename in tqdm(filenames):
            filepath = os.path.join(dirpath, filename)
            file_ext = filename.lower().split(".")[-1]
            if file_ext == "docx":
                text_list, comments = extract_docx(filepath)
                for sent in text_list:
                    for chunk in chunk_sentences(sent):
                        records.append({
                            "File Path": dirpath,
                            "File Name": filename,
                            "File Type": "docx",
                            "Page/Slide/Sheet": None,
                            "Extracted sentence": chunk,
                            "Comments": " | ".join(comments),
                            "Metadata": None
                        })
            elif file_ext == "pptx":
                slides = extract_pptx(filepath)
                for slide_number, slide_text in slides:
                    for chunk in chunk_sentences(slide_text):
                        records.append({
                            "File Path": dirpath,
                            "File Name": filename,
                            "File Type": "pptx",
                            "Page/Slide/Sheet": slide_number,
                            "Extracted sentence": chunk,
                            "Comments": None,
                            "Metadata": None
                        })
            elif file_ext == "xlsx":
                rows = extract_xlsx(filepath)
                for sheet_name, row_idx, row_text in rows:
                    for chunk in chunk_sentences(row_text):
                        records.append({
                            "File Path": dirpath,
                            "File Name": filename,
                            "File Type": "xlsx",
                            "Page/Slide/Sheet": f"{sheet_name}:{row_idx}",
                            "Extracted sentence": chunk,
                            "Comments": None,
                            "Metadata": None
                        })
            elif file_ext == "pdf":
                pages = extract_pdf(filepath, use_ocr=True)
                for page_number, page_text in pages:
                    for chunk in chunk_sentences(page_text):
                        records.append({
                            "File Path": dirpath,
                            "File Name": filename,
                            "File Type": "pdf",
                            "Page/Slide/Sheet": page_number,
                            "Extracted sentence": chunk,
                            "Comments": None,
                            "Metadata": None
                        })
            else:
                continue
    return records

# -------------------------------
# Entry point
# -------------------------------
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Extract text from office documents and PDFs into CSV/XLSX")
    parser.add_argument("root_dir", help="Root directory to walk and extract files from")
    parser.add_argument("--output_csv", default="extracted_corpus.csv", help="Output CSV file")
    parser.add_argument("--output_xlsx", default="extracted_corpus.xlsx", help="Output XLSX file")
    args = parser.parse_args()

    all_records = process_directory(args.root_dir)
    print(f"Total sentences extracted: {len(all_records)}")

    df = pd.DataFrame(all_records)

    # Write CSV
    df.to_csv(args.output_csv, index=False)
    print(f"Saved CSV: {args.output_csv}")

    # Write XLSX
    df.to_excel(args.output_xlsx, index=False)
    print(f"Saved XLSX: {args.output_xlsx}")
